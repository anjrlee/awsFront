import io
import magic
import pytesseract
from pdf2image import convert_from_bytes
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image

def convert_to_text(file) -> (io.BytesIO, str):
    mime = magic.from_buffer(file.read(2048), mime=True)
    file.seek(0)  # reset pointer

    if mime == 'application/pdf':
        return extract_text_from_pdf(file)
    elif mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return extract_text_from_docx(file)
    elif mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return extract_text_from_pptx(file)
    elif mime.startswith('image/'):
        return extract_text_from_image(file)
    else:
        raise ValueError(f"Unsupported file type: {mime}")


def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = "\n".join([page.get_text() for page in doc])
    file_like = io.BytesIO(text.encode("utf-8"))
    return file_like, f"{uuid.uuid4()}.txt"

def extract_text_from_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    file_like = io.BytesIO(text.encode("utf-8"))
    return file_like, f"{uuid.uuid4()}.txt"

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    text = "\n".join(text_runs)
    file_like = io.BytesIO(text.encode("utf-8"))
    return file_like, f"{uuid.uuid4()}.txt"

def extract_text_from_image(file):
    image = Image.open(file)
    text = pytesseract.image_to_string(image)
    file_like = io.BytesIO(text.encode("utf-8"))
    return file_like, f"{uuid.uuid4()}.txt"
